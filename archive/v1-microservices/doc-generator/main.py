import os
from datetime import datetime
from typing import Dict, Any, List
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
import structlog
from jinja2 import Environment, FileSystemLoader
from weasyprint import HTML, CSS
from docx import Document
from pptx import Presentation
from pptx.util import Inches

structlog.configure(
    processors=[structlog.dev.ConsoleRenderer()],
    wrapper_class=structlog.make_filtering_bound_logger(20),
    logger_factory=structlog.PrintLoggerFactory(),
    cache_logger_on_first_use=True
)

logger = structlog.get_logger(__name__)

app = FastAPI(title="Document Generator Service", description="Generate PDF, Word, and PowerPoint documents", version="1.0.0")

app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

class DocumentRequest(BaseModel):
    document_type: str  # pdf, word, powerpoint
    template_name: str = "default"
    data: Dict[str, Any]
    filename: str = None

@app.on_event("startup")
async def startup_event():
    """Initialize document generation service"""
    try:
        # Ensure output directory exists
        os.makedirs("/app/output", exist_ok=True)

        # Create default templates
        create_default_templates()

        logger.info("Document Generator Service started successfully")

    except Exception as e:
        logger.error("Failed to initialize Document Generator service", error=str(e))
        raise

def create_default_templates():
    """Create default templates if they don't exist"""
    os.makedirs("/app/templates", exist_ok=True)

    # Default HTML template for PDF generation
    html_template = """
    <!DOCTYPE html>
    <html>
    <head>
        <title>{{ title | default('Task Report') }}</title>
        <style>
            body { font-family: Arial, sans-serif; margin: 40px; }
            h1 { color: #2c3e50; border-bottom: 2px solid #3498db; }
            h2 { color: #34495e; }
            table { width: 100%; border-collapse: collapse; margin: 20px 0; }
            th, td { border: 1px solid #ddd; padding: 8px; text-align: left; }
            th { background-color: #f2f2f2; }
            .completed { color: #27ae60; }
            .in-progress { color: #f39c12; }
            .not-started { color: #e74c3c; }
        </style>
    </head>
    <body>
        <h1>{{ title | default('Task Report') }}</h1>
        <p><strong>Generated:</strong> {{ generated_date }}</p>

        {% if plan_name %}
        <h2>Plan: {{ plan_name }}</h2>
        {% endif %}

        {% if tasks %}
        <h2>Tasks</h2>
        <table>
            <tr>
                <th>Title</th>
                <th>Status</th>
                <th>Due Date</th>
                <th>Progress</th>
            </tr>
            {% for task in tasks %}
            <tr>
                <td>{{ task.title }}</td>
                <td class="{% if task.percent_complete == 100 %}completed{% elif task.percent_complete > 0 %}in-progress{% else %}not-started{% endif %}">
                    {% if task.percent_complete == 100 %}Completed{% elif task.percent_complete > 0 %}In Progress{% else %}Not Started{% endif %}
                </td>
                <td>{{ task.due_date | default('No due date') }}</td>
                <td>{{ task.percent_complete }}%</td>
            </tr>
            {% endfor %}
        </table>
        {% endif %}

        {% if summary %}
        <h2>Summary</h2>
        <p>{{ summary }}</p>
        {% endif %}
    </body>
    </html>
    """

    with open("/app/templates/default.html", "w") as f:
        f.write(html_template)

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "output_directory": "/app/output",
        "supported_formats": ["pdf", "word", "powerpoint"]
    }

@app.post("/generate")
async def generate_document(request: DocumentRequest):
    """Generate a document based on the request"""
    try:
        # Generate filename if not provided
        if not request.filename:
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            request.filename = f"document_{timestamp}"

        # Add generated date to data
        request.data["generated_date"] = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S")

        if request.document_type.lower() == "pdf":
            return await generate_pdf(request)
        elif request.document_type.lower() == "word":
            return await generate_word(request)
        elif request.document_type.lower() == "powerpoint":
            return await generate_powerpoint(request)
        else:
            raise HTTPException(status_code=400, detail="Unsupported document type")

    except Exception as e:
        logger.error("Error generating document", error=str(e))
        raise HTTPException(status_code=500, detail=f"Document generation failed: {str(e)}")

async def generate_pdf(request: DocumentRequest):
    """Generate PDF document"""
    try:
        # Load Jinja2 template
        env = Environment(loader=FileSystemLoader("/app/templates"))
        template = env.get_template(f"{request.template_name}.html")

        # Render HTML
        html_content = template.render(**request.data)

        # Generate PDF
        output_path = f"/app/output/{request.filename}.pdf"
        HTML(string=html_content).write_pdf(output_path)

        logger.info("PDF generated successfully", filename=request.filename)

        return {
            "success": True,
            "message": "PDF generated successfully",
            "filename": f"{request.filename}.pdf",
            "download_url": f"/download/{request.filename}.pdf"
        }

    except Exception as e:
        logger.error("Error generating PDF", error=str(e))
        raise

async def generate_word(request: DocumentRequest):
    """Generate Word document"""
    try:
        doc = Document()

        # Add title
        title = request.data.get("title", "Task Report")
        doc.add_heading(title, 0)

        # Add generated date
        doc.add_paragraph(f"Generated: {request.data.get('generated_date', '')}")

        # Add plan name if provided
        if "plan_name" in request.data:
            doc.add_heading(f"Plan: {request.data['plan_name']}", level=1)

        # Add tasks table if provided
        if "tasks" in request.data and request.data["tasks"]:
            doc.add_heading("Tasks", level=1)

            table = doc.add_table(rows=1, cols=4)
            table.style = 'Table Grid'

            # Header row
            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = 'Title'
            hdr_cells[1].text = 'Status'
            hdr_cells[2].text = 'Due Date'
            hdr_cells[3].text = 'Progress'

            # Data rows
            for task in request.data["tasks"]:
                row_cells = table.add_row().cells
                row_cells[0].text = str(task.get("title", ""))

                percent = task.get("percent_complete", 0)
                if percent == 100:
                    status = "Completed"
                elif percent > 0:
                    status = "In Progress"
                else:
                    status = "Not Started"
                row_cells[1].text = status

                row_cells[2].text = str(task.get("due_date", "No due date"))
                row_cells[3].text = f"{percent}%"

        # Add summary if provided
        if "summary" in request.data:
            doc.add_heading("Summary", level=1)
            doc.add_paragraph(request.data["summary"])

        # Save document
        output_path = f"/app/output/{request.filename}.docx"
        doc.save(output_path)

        logger.info("Word document generated successfully", filename=request.filename)

        return {
            "success": True,
            "message": "Word document generated successfully",
            "filename": f"{request.filename}.docx",
            "download_url": f"/download/{request.filename}.docx"
        }

    except Exception as e:
        logger.error("Error generating Word document", error=str(e))
        raise

async def generate_powerpoint(request: DocumentRequest):
    """Generate PowerPoint presentation"""
    try:
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = request.data.get("title", "Task Report")
        subtitle.text = f"Generated: {request.data.get('generated_date', '')}"

        # Plan overview slide if plan_name provided
        if "plan_name" in request.data:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = f"Plan: {request.data['plan_name']}"

            tf = body_shape.text_frame
            tf.text = "Plan Overview"

            if "summary" in request.data:
                p = tf.add_paragraph()
                p.text = request.data["summary"]

        # Tasks slide if tasks provided
        if "tasks" in request.data and request.data["tasks"]:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            title_shape.text = "Tasks"

            tf = body_shape.text_frame
            tf.text = "Task Summary"

            completed = sum(1 for task in request.data["tasks"] if task.get("percent_complete", 0) == 100)
            in_progress = sum(1 for task in request.data["tasks"] if 0 < task.get("percent_complete", 0) < 100)
            not_started = sum(1 for task in request.data["tasks"] if task.get("percent_complete", 0) == 0)

            p = tf.add_paragraph()
            p.text = f"Completed: {completed}"

            p = tf.add_paragraph()
            p.text = f"In Progress: {in_progress}"

            p = tf.add_paragraph()
            p.text = f"Not Started: {not_started}"

        # Save presentation
        output_path = f"/app/output/{request.filename}.pptx"
        prs.save(output_path)

        logger.info("PowerPoint presentation generated successfully", filename=request.filename)

        return {
            "success": True,
            "message": "PowerPoint presentation generated successfully",
            "filename": f"{request.filename}.pptx",
            "download_url": f"/download/{request.filename}.pptx"
        }

    except Exception as e:
        logger.error("Error generating PowerPoint presentation", error=str(e))
        raise

@app.get("/download/{filename}")
async def download_file(filename: str):
    """Download generated file"""
    file_path = f"/app/output/{filename}"

    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    return FileResponse(
        path=file_path,
        filename=filename,
        media_type='application/octet-stream'
    )

@app.get("/templates")
async def list_templates():
    """List available templates"""
    templates_dir = "/app/templates"
    if not os.path.exists(templates_dir):
        return {"templates": []}

    templates = [f.replace(".html", "") for f in os.listdir(templates_dir) if f.endswith(".html")]
    return {"templates": templates}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)